"""
document_processor.py — Extraction de texte et découpage en chunks
Moteur Sémantique Multidocument

Supporte :
    - PDF (PyPDF)
    - Word (.docx)
    - PowerPoint (.pptx)
    - Excel (.xlsx)
    - Texte (.txt)
    - RTF (.rtf)
"""

import os
from pypdf import PdfReader
import docx
import pptx
import openpyxl
from striprtf.striprtf import rtf_to_text


def extraire_texte_pdf(pdf_path):
    nom_fichier = os.path.basename(pdf_path)
    reader = PdfReader(pdf_path)
    paragraphes = []
    for num_page, page in enumerate(reader.pages, start=1):
        texte = page.extract_text()
        if texte:
            texte = " ".join(texte.split())
            if texte:
                paragraphes.append({
                    "texte": texte,
                    "page": num_page,
                    "fichier": nom_fichier,
                })
    return paragraphes


def extraire_texte_word(docx_path):
    nom_fichier = os.path.basename(docx_path)
    doc = docx.Document(docx_path)
    paragraphes = []
    for i, para in enumerate(doc.paragraphs, start=1):
        texte = para.text.strip()
        if texte:
            paragraphes.append({
                "texte": texte,
                "page": i,  # Numéro de paragraphe comme page
                "fichier": nom_fichier,
            })
    return paragraphes


def extraire_texte_pptx(pptx_path):
    nom_fichier = os.path.basename(pptx_path)
    prs = pptx.Presentation(pptx_path)
    paragraphes = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        texte = " ".join(slide_text).strip()
        if texte:
            paragraphes.append({
                "texte": texte,
                "page": i,
                "fichier": nom_fichier,
            })
    return paragraphes


def extraire_texte_excel(xlsx_path):
    nom_fichier = os.path.basename(xlsx_path)
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    paragraphes = []
    page_num = 1
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None).strip()
            if row_text:
                paragraphes.append({
                    "texte": row_text,
                    "page": page_num,
                    "fichier": f"{nom_fichier} ({sheet_name})",
                })
                page_num += 1
    return paragraphes


def extraire_texte_txt(txt_path):
    nom_fichier = os.path.basename(txt_path)
    paragraphes = []
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        for i, line in enumerate(f, start=1):
            texte = line.strip()
            if texte:
                paragraphes.append({
                    "texte": texte,
                    "page": i,
                    "fichier": nom_fichier,
                })
    return paragraphes


def extraire_texte_rtf(rtf_path):
    nom_fichier = os.path.basename(rtf_path)
    paragraphes = []
    with open(rtf_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    
    text = rtf_to_text(content)
    for i, line in enumerate(text.split("\n"), start=1):
        texte = line.strip()
        if texte:
            paragraphes.append({
                "texte": texte,
                "page": i,
                "fichier": nom_fichier,
            })
    return paragraphes


def extraire_texte(file_path):
    """
    Détecte automatiquement le type de fichier selon son extension 
    et appelle la bonne fonction.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extraire_texte_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        # NOTE: python-docx ne lit que les .docx, pas les vieux .doc
        # On va essayer de le lire comme docx, si ça plante c'est probablement un vieux doc
        try:
            return extraire_texte_word(file_path)
        except Exception:
            return []
    elif ext in [".pptx", ".ppt"]:
        try:
            return extraire_texte_pptx(file_path)
        except Exception:
            return []
    elif ext in [".xlsx", ".xls"]:
        try:
            return extraire_texte_excel(file_path)
        except Exception:
            return []
    elif ext == ".txt":
        return extraire_texte_txt(file_path)
    elif ext == ".rtf":
        return extraire_texte_rtf(file_path)
    else:
        print(f"⚠️ Avertissement : Type de fichier non supporté : {ext}")
        return []


def decouper_chunks(paragraphes_meta, taille=100, overlap=20):
    """
    Prend une liste de dictionnaires avec métadonnées, les regroupe en chunks
    de maximum 'taille' mots avec un chevauchement de 'overlap' mots.

    Le chevauchement permet de conserver le contexte entre deux chunks consécutifs.
    Exemple : chunk1 = mots 1→100, chunk2 = mots 80→180, chunk3 = mots 160→260

    Chaque chunk conserve le numéro de page du premier mot qui le compose.

    Args:
        paragraphes_meta: Liste de dict {"texte": ..., "page": ..., "fichier": ...}
        taille:           Nombre maximum de mots par chunk (défaut: 100).
        overlap:          Nombre de mots de chevauchement entre chunks (défaut: 20).

    Returns:
        Liste de dict : {"texte": "...", "page": N, "fichier": "nom.pdf"}
    """
    # Construire une liste plate de (mot, page, fichier) pour le sliding window
    mots_avec_meta = []
    for para in paragraphes_meta:
        texte = para["texte"]
        page = para.get("page", 1)
        fichier = para.get("fichier", "")
        for mot in texte.split():
            mots_avec_meta.append((mot, page, fichier))

    if not mots_avec_meta:
        return []

    chunks = []
    pas = max(1, taille - overlap)  # Avancer de (taille - overlap) mots à chaque itération
    i = 0

    while i < len(mots_avec_meta):
        fin = min(i + taille, len(mots_avec_meta))
        fenetre = mots_avec_meta[i:fin]

        texte_chunk = " ".join(m[0] for m in fenetre)
        page_debut = fenetre[0][1]
        fichier = fenetre[0][2]

        chunks.append({
            "texte": texte_chunk,
            "page": page_debut,
            "fichier": fichier,
        })

        # Avancer avec le pas (taille - overlap)
        i += pas

        # Éviter de créer un chunk minuscule à la fin
        if i < len(mots_avec_meta) and len(mots_avec_meta) - i < overlap:
            break

    return chunks

